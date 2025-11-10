import json
import sqlite3
import asyncio
from datetime import datetime
import warnings
from pathlib import Path
from typing import Any, Dict, List

import streamlit as st
import tiktoken

from tools.module2_chunker_tool import (
    DEFAULT_MAX_TOKENS as CHUNK_MAX_TOKENS,
    DEFAULT_OVERLAP,
    adaptive_chunk,
)
from tools.module3_embedder_tool import (
    DEFAULT_COLLECTION,
    DEFAULT_MODEL as EMBEDDING_MODEL,
    DEFAULT_PERSIST_DIR,
    embed_and_store,
    reset_vector_store,
)
from tools.module4_query_rewriter_tool import (
    DEFAULT_MODEL as REWRITER_MODEL,
    _get_client as get_gemini_client,  # type: ignore
)
from tools.module5_context_composer_tool import DEFAULT_MAX_TOKENS as COMPOSER_MAX_TOKENS
from tools.module7_memory_manager_tool import DEFAULT_LOG_PATH, store_feedback
from agents.module6_agentic_rag import (
    DEFAULT_TOP_K,
    handle_query,
    retrieve_relevant_chunks,
)
from google.adk.runners import InMemoryRunner
from agent import root_agent

MODULE1_STAGES = [
    "Chunking",
    "Embedding",
    "Retrieval",
    "Context Assembly",
    "Reasoning",
    "Feedback",
]

ENCODER = tiktoken.get_encoding("cl100k_base")
NOTEBOOK_SUFFIXES = {".ipynb"}
PDF_SUFFIXES = {".pdf"}
PPTX_SUFFIXES = {".pptx"}
BINARY_SKIP_SUFFIXES = {
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".bmp",
    ".svg",
    ".mp4",
}
MAX_FALLBACK_BYTES = 1_000_000  # 1 MB


@st.cache_resource
def _get_agent_runner() -> InMemoryRunner:
    return InMemoryRunner(agent=root_agent)


def _format_agent_events(events: list) -> str:
    """Render ADK agent events into markdown for the chat panel."""

    outputs: list[str] = []
    for event in events:
        if event.author != root_agent.name:
            continue
        if not event.content or not event.content.parts:
            continue
        for part in event.content.parts:
            if getattr(part, "text", None):
                outputs.append(part.text)
            elif getattr(part, "function_call", None):
                payload = json.dumps(part.function_call.args or {}, indent=2)
                outputs.append(
                    f"**Tool Call** `{part.function_call.name}`\n```\n{payload}\n```"
                )
            elif getattr(part, "function_response", None):
                payload = json.dumps(part.function_response.response, indent=2)
                outputs.append(
                    f"**Tool Response** `{part.function_response.name}`\n```\n{payload}\n```"
                )
    return "\n\n".join(outputs).strip() or "_Agent produced no textual response._"


def _format_timestamp(epoch_seconds: float) -> str:
    return datetime.fromtimestamp(epoch_seconds).strftime("%d %b %Y · %H:%M")


def _get_vector_store_stats(persist_dir: Path) -> Dict[str, Any]:
    stats: Dict[str, Any] = {
        "indexed_vectors": 0,
        "collections": 0,
        "last_updated": "Not indexed",
    }

    db_path = persist_dir / "chroma.sqlite3"
    if not db_path.exists():
        return stats

    conn: sqlite3.Connection | None = None
    try:
        conn = sqlite3.connect(db_path)
        cur = conn.cursor()
        stats["indexed_vectors"] = int(cur.execute("SELECT COUNT(*) FROM embeddings").fetchone()[0])
        try:
            stats["collections"] = int(cur.execute("SELECT COUNT(*) FROM collections").fetchone()[0])
        except sqlite3.OperationalError:
            stats["collections"] = 1 if stats["indexed_vectors"] else 0
    except sqlite3.Error:
        return stats
    finally:
        if conn is not None:
            conn.close()

    try:
        stats["last_updated"] = _format_timestamp(db_path.stat().st_mtime)
    except OSError:
        pass
    return stats


def _render_status_overview(stats: Dict[str, Any]) -> None:
    st.markdown("### System Snapshot")
    col1, col2, col3 = st.columns(3)
    col1.metric("Indexed Vectors", f"{stats['indexed_vectors']:,}")
    col2.metric("Collections", stats["collections"])
    col3.metric("Last Updated", stats["last_updated"])
    st.caption(
        "Vector counts update after a successful ingest run. Use the sidebar to rebuild embeddings."
    )


def _render_memory_summary(history: List[dict]) -> None:
    if not history:
        st.info("No feedback captured yet. Run a query and submit feedback to populate this view.")
        return

    st.markdown("#### Feedback Timeline")
    positive = sum(1 for entry in history if entry.get("feedback") == "positive")
    negative = sum(1 for entry in history if entry.get("feedback") == "negative")
    pending = sum(1 for entry in history if entry.get("feedback") == "pending")
    col1, col2, col3 = st.columns(3)
    col1.metric("Positive", positive)
    col2.metric("Pending", pending)
    col3.metric("Negative", negative)

    st.markdown("#### Recent Entries")
    for entry in reversed(history[-20:]):
        with st.container():
            timestamp = entry.get("timestamp", "Unknown time")
            feedback = entry.get("feedback", "pending").title()
            st.markdown(
                f"**{feedback}** · `{timestamp}`\n\n"
                f"**Query:** {entry.get('query', 'N/A')}"
            )
            chunks = entry.get("retrieved_chunks") or []
            if chunks:
                with st.expander("Retrieved Chunks", expanded=False):
                    for idx, chunk in enumerate(chunks, start=1):
                        st.markdown(f"**Chunk {idx}:** {chunk}")
            notes = entry.get("notes")
            if notes:
                st.caption(f"Notes: {notes}")
            st.divider()


def _load_feedback_log(path: Path) -> list[dict]:
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text("utf-8"))[-10:]
    except json.JSONDecodeError:
        return []


def _render_chunk_list(chunks: list[str]) -> None:
    if not chunks:
        st.info("No chunks retrieved. Ingest documents via Module 3 to populate Chroma.")
        return
    for idx, chunk in enumerate(chunks, start=1):
        with st.expander(f"Chunk #{idx}"):
            st.write(chunk)
            st.caption(f"Token length: {len(ENCODER.encode(chunk))}")


def _render_feedback_panel(result: dict[str, str]) -> None:
    st.subheader("Module 7 · Feedback & Memory")
    question = result["question"]
    base_key = f"feedback_{abs(hash(question))}"
    label_key = f"{base_key}_label"
    notes_key = f"{base_key}_notes"
    button_key = f"{base_key}_submit"

    if label_key not in st.session_state:
        st.session_state[label_key] = "pending"
    if notes_key not in st.session_state:
        st.session_state[notes_key] = ""
    reset_key = f"{notes_key}_reset"
    if st.session_state.pop(reset_key, False):
        st.session_state[notes_key] = ""

    options = ["positive", "pending", "negative"]
    feedback = st.selectbox("Feedback label", options, key=label_key)
    notes = st.text_area(
        "Notes",
        placeholder="Optionally add comments for retraining…",
        key=notes_key,
    )
    if st.button("Store Feedback", key=button_key):
        msg = store_feedback(
            query=result["question"],
            retrieved_chunks=result.get("chunks", []),
            feedback=feedback,
            notes=notes or None,
        )
        st.success(msg)
        st.session_state[reset_key] = True

    history = _load_feedback_log(DEFAULT_LOG_PATH)
    if history:
        with st.expander("Recent Feedback Entries", expanded=False):
            for entry in reversed(history):
                st.json(entry, expanded=False)


def _attempt_collection_preview(top_k: int) -> None:
    st.subheader("Module 3 · Embeddings & Vector Store")
    info_cols = st.columns(3)
    info_cols[0].metric("Embedding Model", EMBEDDING_MODEL)
    info_cols[1].metric("Collection", DEFAULT_COLLECTION)
    info_cols[2].metric("Top-K", top_k)

    try:
        client = sqlite3.connect(Path(DEFAULT_PERSIST_DIR) / "chroma.sqlite3")
        cur = client.cursor()
        count = cur.execute("SELECT COUNT(*) FROM embeddings").fetchone()[0]
        st.caption(f"Indexed embeddings: {count}")
    except (sqlite3.Error, FileNotFoundError):
        st.caption("Indexed embeddings: unavailable (ingest documents to enable retrieval).")
    finally:
        try:
            client.close()
        except Exception:
            pass


def _load_file_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in NOTEBOOK_SUFFIXES:
        try:
            data = json.loads(path.read_text("utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise RuntimeError("Failed to parse notebook JSON") from exc
        cells = data.get("cells", [])
        assembled: list[str] = []
        for cell in cells:
            source = "".join(cell.get("source", []))
            if source.strip():
                assembled.append(source)
        return "\n".join(assembled)
    if suffix in PPTX_SUFFIXES:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:  # pragma: no cover - optional dependency
            raise RuntimeError(
                "PowerPoint support requires the 'python-pptx' package. Install it to ingest .pptx files."
            ) from exc
        presentation = Presentation(str(path))
        slide_text: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
        return "\n".join(slide_text)
    if suffix in PDF_SUFFIXES:
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError as exc:  # pragma: no cover - optional dependency
            raise RuntimeError(
                "PDF support requires the 'pypdf' package. Install it to ingest PDFs."
            ) from exc
        try:
            reader = PdfReader(str(path))
        except Exception as exc:
            raise RuntimeError(f"Failed to parse PDF: {exc}") from exc
        text_parts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(text_parts)

    if suffix in BINARY_SKIP_SUFFIXES:
        return ""

    sample_size = min(MAX_FALLBACK_BYTES, path.stat().st_size)
    if sample_size <= 0:
        return ""
    with path.open("rb") as binary_file:
        sample = binary_file.read(min(8192, sample_size))
    if b"\x00" in sample:
        warnings.warn(f"{path.name}: binary file detected; skipped.", RuntimeWarning)
        return ""
    try:
        return path.read_text("utf-8")
    except UnicodeDecodeError:
        try:
            return path.read_text("utf-8", errors="ignore")
        except OSError as exc:
            raise RuntimeError(f"Failed to read file: {exc}") from exc
    except OSError as exc:
        raise RuntimeError(f"Failed to read file: {exc}") from exc


def _collect_documents(source_dir: Path) -> list[tuple[Path, str]]:
    docs: list[tuple[Path, str]] = []
    for file_path in sorted(source_dir.rglob("*")):
        if not file_path.is_file():
            continue
        if file_path.suffix.lower() in BINARY_SKIP_SUFFIXES:
            st.info(f"{file_path.name}: binary asset skipped.")
            continue
        try:
            content = _load_file_text(file_path)
        except RuntimeError as exc:
            st.info(f"{file_path.name}: {exc}")
            continue
        if not content.strip():
            st.info(f"{file_path.name}: empty or unreadable, skipped.")
            continue
        docs.append((file_path, content))
        if len(docs) >= st.session_state.get("_ingest_file_limit", float("inf")):
            break
    return docs


def _ingest_directory(
    directory: str,
    *,
    chunk_tokens: int = CHUNK_MAX_TOKENS,
    chunk_overlap: int = DEFAULT_OVERLAP,
    embedding_model: str = EMBEDDING_MODEL,
    collection_name: str = DEFAULT_COLLECTION,
    persist_directory: Path = DEFAULT_PERSIST_DIR,
    reset_store: bool = True,
) -> str:
    source_dir = Path(directory).expanduser()
    if not source_dir.exists() or not source_dir.is_dir():
        raise FileNotFoundError(f"Directory not found: {source_dir}")

    documents = _collect_documents(source_dir)
    if not documents:
        raise RuntimeError("No readable documents found in the provided directory.")

    if reset_store:
        reset_vector_store(
            collection_name=collection_name,
            persist_directory=persist_directory,
            drop_storage=True,
        )

    chunks: list[str] = []
    metadata: list[dict] = []
    for file_path, content in documents:
        file_chunks = adaptive_chunk(
            content,
            max_tokens=chunk_tokens,
            overlap=chunk_overlap,
        )
        for idx, chunk in enumerate(file_chunks):
            if not chunk.strip():
                continue
            chunks.append(chunk)
            metadata.append(
                {
                    "source_path": str(file_path),
                    "file_name": file_path.name,
                    "chunk_index": idx,
                }
            )

    if not chunks:
        raise RuntimeError("Chunking produced no content; check the source files.")

    response = embed_and_store(
        chunks,
        metadata_list=metadata,
        model_name=embedding_model,
        collection_name=collection_name,
        persist_directory=persist_directory,
        id_prefix="ingested",
    )
    return f"{response} (from {len(documents)} documents)"


def main() -> None:
    st.set_page_config(page_title="Agentic RAG Console", layout="wide")
    st.title("Agentic RAG Pipeline")
    st.markdown(
        "This Streamlit console wraps the Module 1-7 pipeline to visualise each stage."
    )

    with st.sidebar:
        st.header("Pipeline Settings")
        top_k = st.slider("Top-K Retrieval", 1, 10, DEFAULT_TOP_K)
        max_context_tokens = st.slider("Context Token Budget", 500, 6000, COMPOSER_MAX_TOKENS, step=100)
        st.write("\nGemini model:", REWRITER_MODEL)
        st.write("Chroma directory:", str(DEFAULT_PERSIST_DIR))

        st.divider()
        st.subheader("Module 3 · Rebuild Embeddings")
        folder_path = st.text_input(
            "Source folder",
            value=str(Path.cwd() / "data" / "sample_docs"),
            help="Provide a directory containing text (.txt/.md/.json) or PDF files.",
        )
        reset_collection = st.checkbox(
            "Reset vector store before ingest",
            value=True,
            help="Clears existing embeddings before loading new documents.",
        )
        ingest_status = st.empty()
        if st.button("Rebuild Vector Store", type="secondary"):
            with ingest_status, st.spinner("Processing documents..."):
                try:
                    message = _ingest_directory(
                        folder_path,
                        chunk_tokens=CHUNK_MAX_TOKENS,
                        chunk_overlap=DEFAULT_OVERLAP,
                        embedding_model=EMBEDDING_MODEL,
                        collection_name=DEFAULT_COLLECTION,
                        persist_directory=Path(DEFAULT_PERSIST_DIR),
                        reset_store=reset_collection,
                    )
                except Exception as exc:
                    st.error(f"Ingestion failed: {exc}")
                else:
                    st.success(message)

    stats = _get_vector_store_stats(Path(DEFAULT_PERSIST_DIR))
    _render_status_overview(stats)
    st.divider()

    console_tab, vector_tab, memory_tab, chat_tab = st.tabs(
        ["Pipeline Console", "Vector Store", "Feedback Memory", "Agent Chat"]
    )

    with console_tab:
        question = st.text_input(
            "Ask a question", placeholder="What is the policy for remote work?"
        )
        submit = st.button("Run Agentic RAG", type="primary")

        with st.expander("Module 1 · Context Engineering Flow", expanded=False):
            st.write(" → ".join(MODULE1_STAGES))

        col_chunk, col_context = st.columns(2)
        with col_chunk:
            st.subheader("Module 2 · Chunking")
            st.caption(
                f"Adaptive chunking defaults — max_tokens={CHUNK_MAX_TOKENS}, overlap={DEFAULT_OVERLAP}"
            )
        with col_context:
            st.subheader("Module 5 · Context Composer")
            st.caption(f"Token budget ≤ {max_context_tokens}")

        result: Dict[str, Any] | None = None
        if submit and question:
            try:
                result = handle_query(
                    question,
                    top_k=top_k,
                    max_context_tokens=max_context_tokens,
                )
            except Exception as exc:
                st.error(f"Pipeline failure: {exc}")
                return
            st.session_state["__last_result"] = result
        else:
            cached = st.session_state.get("__last_result")
            if isinstance(cached, dict):
                result = cached

        if result:
            st.subheader("Module 4 · Query Understanding")
            st.code(result["rewritten_query"], language="markdown")

            st.subheader("Module 3 · Retrieval Output")
            _render_chunk_list(result.get("chunks", []))

            st.subheader("Module 5 · Composed Context")
            if context := result.get("context"):
                st.text_area("Context", context, height=240)
                st.caption(f"Final token usage: {len(ENCODER.encode(context))}")
            else:
                st.warning("Context is empty. Ensure embeddings are ingested into Chroma.")

            st.subheader("Module 6 · Reasoning")
            st.markdown(result.get("answer", "No answer."))

            _render_feedback_panel(result)

            st.info(
                "Gemini client re-used from Module 4; make sure `.env` contains either "
                "`GOOGLE_AI_API_KEY` or Vertex AI credentials."
            )

    with vector_tab:
        st.subheader("Vector Store Health")
        _attempt_collection_preview(top_k)
        st.markdown(
            "Use the sidebar controls to rebuild embeddings from your document directory. "
            "The metrics above update once ingestion completes."
        )

    with memory_tab:
        history = _load_feedback_log(DEFAULT_LOG_PATH)
        _render_memory_summary(history)

    with chat_tab:
        st.subheader("Ask the Agent")
        st.caption(
            "Chat with the Agentic RAG assistant. The agent runs the full pipeline on each question "
            "and cites the supporting context when available."
        )

        chat_state = st.session_state.setdefault("agent_chat_history", [])
        for message in chat_state:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

        prompt = st.chat_input("Type your question")
        if prompt:
            chat_state.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)

            with st.chat_message("assistant"):
                with st.spinner("Consulting the vector store..."):
                    runner = _get_agent_runner()
                    try:
                        events = asyncio.run(
                            runner.run_debug(
                                prompt,
                                user_id="streamlit_user",
                                session_id="streamlit_chat",
                                quiet=True,
                            )
                        )
                    except Exception as exc:  # pragma: no cover - runtime feedback
                        error_msg = f"Agent execution failed: {exc}"
                        st.error(error_msg)
                        chat_state.append({"role": "assistant", "content": error_msg})
                    else:
                        agent_reply = _format_agent_events(events)
                        st.markdown(agent_reply)
                        chat_state.append({"role": "assistant", "content": agent_reply})

        if st.button("Clear conversation", key="clear_agent_chat"):
            st.session_state["agent_chat_history"] = []
            _get_agent_runner.clear()
            if hasattr(st, "rerun"):
                st.rerun()
            elif hasattr(st, "experimental_rerun"):  # pragma: no cover - legacy fallback
                st.experimental_rerun()

if __name__ == "__main__":
    main()
